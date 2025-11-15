#!/usr/bin/env python3
"""
NUS Canvas Deadline Scraper
Automatically extracts assignment deadlines and exam dates from Canvas course materials.
"""

import os
import json
import logging
import argparse
import re
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional
import requests
from dotenv import load_dotenv
import anthropic

# Import document parsers
try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logging.warning("PyPDF2 not available - PDF parsing disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available - DOCX parsing disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PPTX parsing disabled")

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('scraper.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configuration
CANVAS_API_TOKEN = os.getenv('CANVAS_API_TOKEN')
CANVAS_API_URL = os.getenv('CANVAS_API_URL', 'https://canvas.nus.edu.sg/api/v1')
ANTHROPIC_API_KEY = os.getenv('ANTHROPIC_API_KEY')

# Keywords to identify intro/syllabus documents
INTRO_KEYWORDS = [
    # Basic intro keywords
    'intro', 'introduction', 'syllabus', 'course outline',
    'schedule', 'overview', 'course info', 'course information',
    'module information', 'module info',

    # Schedule and assessment keywords
    'course schedule', 'assessment', 'assessments',

    # Week/Lecture variations
    'week 1', 'week1', 'week 0', 'week0',
    'lecture 1', 'lecture1', 'lecture 0', 'lecture0',
    'lec 1', 'lec1', 'lec 0', 'lec0',

    # Topic variations
    'topic 0', 'topic_0', 'topic0',

    # Year/semester indicators combined with module codes
    'ay24', 'ay25', 'ay26', 's1', 's2',
]

# Regex patterns for intro documents
INTRO_PATTERNS = [
    # Lecture number patterns: L0, L00, L01, L1, L01a, etc.
    r'[_\s\-]?l0+[a-z]?[_\s\-]',  # L0, L00, L0a, L00a
    r'[_\s\-]?l0*1[a-z]?[_\s\-]',  # L1, L01, L001, L1a, L01a

    # Lec/Lecture patterns
    r'lec[_\s\-]?0+[_\s\-]',  # lec0, lec00, lec_0
    r'lec[_\s\-]?0*1[_\s\-]',  # lec1, lec01, lec_1

    # Topic patterns
    r'topic[_\s\-]?0',  # topic0, topic_0, topic 0

    # Week patterns
    r'week[_\s\-]?0*1[_\s\-]',  # week1, week01, week_1
]

# Directory for downloaded files
DOWNLOAD_DIR = Path('downloaded_materials')
DOWNLOAD_DIR.mkdir(exist_ok=True)


class CanvasClient:
    """Client for Canvas LMS API"""

    def __init__(self, api_url: str, api_token: str):
        self.api_url = api_url.rstrip('/')
        self.headers = {
            'Authorization': f'Bearer {api_token}',
            'Accept': 'application/json'
        }

    def _get(self, endpoint: str, params: Optional[Dict] = None) -> Any:
        """Make GET request to Canvas API"""
        url = f"{self.api_url}/{endpoint.lstrip('/')}"
        try:
            response = requests.get(url, headers=self.headers, params=params)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            logger.error(f"API request failed: {e}")
            return None

    def get_current_user(self) -> Dict:
        """Get current user information"""
        return self._get('/users/self')

    def get_enrolled_courses(self) -> List[Dict]:
        """Get all enrolled courses for current user"""
        courses = self._get('/courses', params={
            'enrollment_state': 'active',
            'per_page': 100
        })
        return courses if courses else []

    def get_course_files(self, course_id: int) -> List[Dict]:
        """Get all files for a course"""
        files = self._get(f'/courses/{course_id}/files', params={
            'per_page': 100
        })
        return files if files else []

    def get_course_announcements(self, course_id: int, limit: int = 50) -> List[Dict]:
        """Get recent announcements for a course"""
        announcements = self._get(f'/courses/{course_id}/discussion_topics', params={
            'only_announcements': True,
            'per_page': limit,
            'order_by': 'recent_activity'
        })
        return announcements if announcements else []

    def download_file(self, file_url: str, save_path: Path) -> bool:
        """Download a file from Canvas"""
        try:
            response = requests.get(file_url, headers=self.headers, stream=True)
            response.raise_for_status()

            with open(save_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            logger.info(f"Downloaded: {save_path.name}")
            return True
        except Exception as e:
            logger.error(f"Failed to download {save_path.name}: {e}")
            return False


class DocumentParser:
    """Parse different document formats to extract text"""

    @staticmethod
    def parse_pdf(file_path: Path) -> str:
        """Extract text from PDF"""
        if not PYPDF2_AVAILABLE:
            return ""

        try:
            reader = PdfReader(str(file_path))
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse PDF {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_docx(file_path: Path) -> str:
        """Extract text from DOCX"""
        if not DOCX_AVAILABLE:
            return ""

        try:
            doc = DocxDocument(str(file_path))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Failed to parse DOCX {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint"""
        if not PPTX_AVAILABLE:
            return ""

        try:
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_txt(file_path: Path) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to parse TXT {file_path.name}: {e}")
            return ""

    @classmethod
    def parse_document(cls, file_path: Path) -> str:
        """Parse document based on file extension"""
        ext = file_path.suffix.lower()

        parsers = {
            '.pdf': cls.parse_pdf,
            '.docx': cls.parse_docx,
            '.doc': cls.parse_docx,
            '.pptx': cls.parse_pptx,
            '.ppt': cls.parse_pptx,
            '.txt': cls.parse_txt,
        }

        parser = parsers.get(ext)
        if parser:
            return parser(file_path)
        else:
            logger.warning(f"Unsupported file format: {ext}")
            return ""


class DeadlineExtractor:
    """Extract deadlines from documents using Claude AI"""

    def __init__(self, api_key: str):
        self.client = anthropic.Anthropic(api_key=api_key)

    def extract_deadlines(self, text: str, course_name: str, current_year: int) -> List[Dict]:
        """Use Claude to extract deadlines from document text"""

        if not text.strip():
            logger.warning(f"No text to extract from {course_name}")
            return []

        # Truncate text if too long (Claude has token limits)
        max_chars = 100000
        if len(text) > max_chars:
            text = text[:max_chars] + "\n... [truncated]"

        prompt = f"""You are analyzing a course document for "{course_name}".
Extract ALL assignment deadlines, exam dates, quiz dates, project deadlines, and any other important academic dates.

Current year context: {current_year} (use this for year inference if not explicitly stated)

Document text:
{text}

Please extract and return ONLY a JSON array of deadline objects. Each object should have:
- "date": ISO format date (YYYY-MM-DD) - infer year from context if not stated
- "title": Brief description of the deadline
- "type": One of ["assignment", "exam", "quiz", "project", "presentation", "other"]
- "weight": Percentage weight if mentioned (or null)
- "notes": Any additional relevant information

Return ONLY the JSON array, no other text. If no deadlines found, return an empty array [].

Example format:
[
  {{"date": "2024-09-15", "title": "Assignment 1", "type": "assignment", "weight": 10, "notes": "Submit via Canvas"}},
  {{"date": "2024-11-20", "title": "Final Exam", "type": "exam", "weight": 50, "notes": "Open book"}}
]
"""

        try:
            logger.info(f"Extracting deadlines from {course_name} using Claude...")

            message = self.client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=4096,
                temperature=0,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )

            response_text = message.content[0].text.strip()

            # Try to parse JSON response
            # Remove markdown code blocks if present
            if response_text.startswith('```'):
                response_text = response_text.split('```')[1]
                if response_text.startswith('json'):
                    response_text = response_text[4:]
                response_text = response_text.strip()

            deadlines = json.loads(response_text)

            # Add course name to each deadline
            for deadline in deadlines:
                deadline['course'] = course_name

            logger.info(f"Extracted {len(deadlines)} deadlines from {course_name}")
            return deadlines

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON response for {course_name}: {e}")
            logger.debug(f"Response was: {response_text}")
            return []
        except Exception as e:
            logger.error(f"Failed to extract deadlines for {course_name}: {e}")
            return []

    def extract_from_announcements(self, announcements: List[Dict], course_name: str, current_year: int) -> List[Dict]:
        """Extract deadline updates from Canvas announcements"""
        all_updates = []

        for announcement in announcements:
            title = announcement.get('title', '')
            message = announcement.get('message', '')
            posted_at = announcement.get('posted_at', '')

            # Combine title and message
            full_text = f"Title: {title}\n\n{message}"

            # Skip if no deadline-related keywords
            deadline_keywords = ['deadline', 'due', 'extended', 'postponed', 'rescheduled',
                                'changed', 'updated', 'new date', 'assignment', 'exam', 'quiz']
            if not any(keyword in full_text.lower() for keyword in deadline_keywords):
                continue

            logger.info(f"  Checking announcement: {title}")

            # Use AI to extract deadline updates
            prompt = f"""You are analyzing a Canvas announcement from "{course_name}" posted on {posted_at}.
Extract any deadline changes, updates, or new deadlines mentioned.

Current year context: {current_year}

Announcement:
{full_text[:5000]}

Please extract and return ONLY a JSON array of deadline objects. Each object should have:
- "date": ISO format date (YYYY-MM-DD)
- "title": Brief description
- "type": One of ["assignment", "exam", "quiz", "project", "presentation", "other"]
- "weight": Percentage weight if mentioned (or null)
- "notes": Any additional info, especially if this is a CHANGE/UPDATE
- "is_update": true if this modifies an existing deadline, false if new

Return ONLY the JSON array, no other text. If no deadlines found, return [].

Example:
[
  {{"date": "2024-10-01", "title": "Assignment 2", "type": "assignment", "weight": 15,
    "notes": "EXTENDED from Sep 24 to Oct 1", "is_update": true}}
]
"""

            try:
                message = self.client.messages.create(
                    model="claude-sonnet-4-20250514",
                    max_tokens=2048,
                    temperature=0,
                    messages=[{"role": "user", "content": prompt}]
                )

                response_text = message.content[0].text.strip()

                # Remove markdown code blocks
                if response_text.startswith('```'):
                    response_text = response_text.split('```')[1]
                    if response_text.startswith('json'):
                        response_text = response_text[4:]
                    response_text = response_text.strip()

                updates = json.loads(response_text)

                for update in updates:
                    update['course'] = course_name
                    update['source'] = f"Announcement: {title}"
                    update['announcement_date'] = posted_at

                all_updates.extend(updates)
                logger.info(f"    Found {len(updates)} deadline updates")

            except json.JSONDecodeError:
                logger.debug(f"No valid JSON from announcement: {title}")
            except Exception as e:
                logger.debug(f"Error processing announcement: {e}")

        return all_updates


def is_intro_document(filename: str) -> bool:
    """Check if filename suggests it's an intro/syllabus document"""
    filename_lower = filename.lower()

    # Check keywords first
    if any(keyword in filename_lower for keyword in INTRO_KEYWORDS):
        return True

    # Check regex patterns
    for pattern in INTRO_PATTERNS:
        if re.search(pattern, filename_lower):
            return True

    return False


def detect_changes(old_deadlines: List[Dict], new_deadlines: List[Dict]) -> Dict[str, List[Dict]]:
    """Detect changes between old and new deadline lists"""
    changes = {
        'added': [],
        'removed': [],
        'modified': []
    }

    # Create lookup dictionaries (key: course + title)
    def make_key(d):
        return f"{d.get('course', '')}_{d.get('title', '')}".lower()

    old_dict = {make_key(d): d for d in old_deadlines}
    new_dict = {make_key(d): d for d in new_deadlines}

    # Find added and modified
    for key, new_item in new_dict.items():
        if key not in old_dict:
            changes['added'].append(new_item)
        else:
            old_item = old_dict[key]
            # Check if date changed
            if old_item.get('date') != new_item.get('date'):
                change_info = new_item.copy()
                change_info['old_date'] = old_item.get('date')
                changes['modified'].append(change_info)

    # Find removed
    for key, old_item in old_dict.items():
        if key not in new_dict:
            changes['removed'].append(old_item)

    return changes


def format_deadlines_markdown(deadlines: List[Dict], changes: Optional[Dict] = None) -> str:
    """Format deadlines as markdown"""
    if not deadlines:
        return "No deadlines found."

    # Sort by date
    sorted_deadlines = sorted(deadlines, key=lambda x: x.get('date', '9999-12-31'))

    md = f"# Assignment Deadlines - Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"

    # Add changes summary if provided
    if changes and (changes['added'] or changes['removed'] or changes['modified']):
        md += "## ⚠️ RECENT CHANGES DETECTED\n\n"

        if changes['modified']:
            md += "### 📅 Date Changes:\n"
            for item in changes['modified']:
                old_date = item.get('old_date', 'Unknown')
                new_date = item.get('date', 'Unknown')
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title}: ~~{old_date}~~ → **{new_date}**\n"
            md += "\n"

        if changes['added']:
            md += "### ✨ New Deadlines:\n"
            for item in changes['added']:
                date = item.get('date', 'TBD')
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title} ({date})\n"
            md += "\n"

        if changes['removed']:
            md += "### 🗑️ Removed/Cancelled:\n"
            for item in changes['removed']:
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title}\n"
            md += "\n"

        md += "---\n\n"

    current_month = None
    for deadline in sorted_deadlines:
        date_str = deadline.get('date', 'TBD')

        # Add month header
        if date_str != 'TBD':
            try:
                date_obj = datetime.fromisoformat(date_str)
                month_header = date_obj.strftime('%B %Y')
                if month_header != current_month:
                    current_month = month_header
                    md += f"\n## {month_header}\n\n"

                formatted_date = date_obj.strftime('%a, %d %b')
            except:
                formatted_date = date_str
        else:
            formatted_date = 'TBD'
            if current_month != 'TBD':
                current_month = 'TBD'
                md += f"\n## Date TBD\n\n"

        # Add deadline entry
        title = deadline.get('title', 'Untitled')
        course = deadline.get('course', 'Unknown Course')
        weight = deadline.get('weight')
        weight_str = f" ({weight}%)" if weight else ""
        deadline_type = deadline.get('type', 'other').upper()
        notes = deadline.get('notes', '')
        is_update = deadline.get('is_update', False)
        source = deadline.get('source', '')

        md += f"### {formatted_date} - **{course}**\n"
        md += f"- **[{deadline_type}]** {title}{weight_str}\n"

        # Highlight updates from announcements
        if is_update or 'EXTENDED' in notes.upper() or 'CHANGED' in notes.upper():
            md += f"  - 🔔 **UPDATE:** {notes}\n"
        elif notes:
            md += f"  - _{notes}_\n"

        if source:
            md += f"  - _Source: {source}_\n"

        md += "\n"

    return md


def main():
    """Main execution function"""
    parser = argparse.ArgumentParser(description='NUS Canvas Deadline Scraper')
    parser.add_argument('--display', action='store_true', help='Display deadlines in terminal')
    parser.add_argument('--search', type=str, help='Search keywords (comma-separated)')
    args = parser.parse_args()

    # Validate configuration
    if not CANVAS_API_TOKEN:
        logger.error("CANVAS_API_TOKEN not set in .env file")
        return

    if not ANTHROPIC_API_KEY:
        logger.error("ANTHROPIC_API_KEY not set in .env file")
        return

    logger.info("Starting Canvas Deadline Scraper...")

    # Load previous deadlines for change detection
    old_deadlines = []
    if Path('deadlines.json').exists():
        try:
            with open('deadlines.json', 'r', encoding='utf-8') as f:
                old_deadlines = json.load(f)
            logger.info(f"Loaded {len(old_deadlines)} previous deadlines for change detection")
        except Exception as e:
            logger.warning(f"Could not load previous deadlines: {e}")

    # Initialize clients
    canvas = CanvasClient(CANVAS_API_URL, CANVAS_API_TOKEN)
    extractor = DeadlineExtractor(ANTHROPIC_API_KEY)

    # Get current user
    user = canvas.get_current_user()
    if not user:
        logger.error("Failed to authenticate with Canvas API")
        return

    logger.info(f"Authenticated as: {user.get('name', 'Unknown')}")

    # Get enrolled courses
    courses = canvas.get_enrolled_courses()
    logger.info(f"Found {len(courses)} enrolled courses")

    all_deadlines = []
    current_year = datetime.now().year

    # Process each course
    for course in courses:
        course_id = course.get('id')
        course_name = course.get('name', f'Course {course_id}')

        logger.info(f"\nProcessing: {course_name}")

        # Get course files
        files = canvas.get_course_files(course_id)
        logger.info(f"  Found {len(files)} files")

        # Filter for intro documents
        intro_files = [f for f in files if is_intro_document(f.get('filename', ''))]

        if not intro_files:
            logger.warning(f"  No intro documents found for {course_name}")
            logger.info(f"  Available files (first 10):")
            for idx, f in enumerate(files[:10]):
                logger.info(f"    - {f.get('filename', 'unnamed')}")
            if len(files) > 10:
                logger.info(f"    ... and {len(files) - 10} more")
            logger.info(f"  Tip: Intro docs should contain keywords like: intro, syllabus, L0, L01, Topic 0, etc.")
            continue

        logger.info(f"  Found {len(intro_files)} intro documents")

        # Process each intro document
        for file_info in intro_files:
            filename = file_info.get('filename', 'unnamed')
            file_url = file_info.get('url')

            if not file_url:
                continue

            # Download file
            course_dir = DOWNLOAD_DIR / course_name.replace('/', '_')
            course_dir.mkdir(exist_ok=True)
            save_path = course_dir / filename

            if not save_path.exists():
                if not canvas.download_file(file_url, save_path):
                    continue

            # Parse document
            logger.info(f"  Parsing: {filename}")
            text = DocumentParser.parse_document(save_path)

            if not text.strip():
                logger.warning(f"  No text extracted from {filename}")
                continue

            # Extract deadlines
            deadlines = extractor.extract_deadlines(text, course_name, current_year)
            all_deadlines.extend(deadlines)

        # Check announcements for deadline updates
        logger.info(f"  Checking announcements for updates...")
        announcements = canvas.get_course_announcements(course_id, limit=50)
        logger.info(f"  Found {len(announcements)} announcements")

        if announcements:
            updates = extractor.extract_from_announcements(announcements, course_name, current_year)
            if updates:
                logger.info(f"  Extracted {len(updates)} deadline updates from announcements")
                all_deadlines.extend(updates)

    # Detect changes
    changes = None
    if old_deadlines:
        logger.info(f"\nDetecting changes from previous run...")
        changes = detect_changes(old_deadlines, all_deadlines)
        total_changes = len(changes['added']) + len(changes['removed']) + len(changes['modified'])
        if total_changes > 0:
            logger.info(f"  🔔 Found {total_changes} changes:")
            logger.info(f"    - {len(changes['added'])} added")
            logger.info(f"    - {len(changes['modified'])} date changes")
            logger.info(f"    - {len(changes['removed'])} removed")
        else:
            logger.info(f"  No changes detected")

    # Save results
    logger.info(f"\nTotal deadlines extracted: {len(all_deadlines)}")

    # Save JSON
    with open('deadlines.json', 'w', encoding='utf-8') as f:
        json.dump(all_deadlines, f, indent=2, ensure_ascii=False)
    logger.info("Saved to: deadlines.json")

    # Save Markdown (with changes highlighted)
    markdown_content = format_deadlines_markdown(all_deadlines, changes)
    with open('deadlines.md', 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    logger.info("Saved to: deadlines.md")

    # Save changes summary if any
    if changes and (changes['added'] or changes['removed'] or changes['modified']):
        with open('changes.json', 'w', encoding='utf-8') as f:
            json.dump(changes, f, indent=2, ensure_ascii=False)
        logger.info("Saved changes to: changes.json")

    # Display if requested
    if args.display:
        print("\n" + "="*80)
        print(markdown_content)
        print("="*80)

    logger.info("\nDone! Check deadlines.md for formatted output.")


if __name__ == "__main__":
    main()
